import os
from io import BytesIO
from datetime import datetime, date
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from sqlalchemy import text
from sqlalchemy.orm import Session

# --- MODEL IMPORTS ---
from app.models.tank_header import Tank
from app.models.tank_details import TankDetails
from app.models.tank_images import TankImage
from app.models.tank_regulations import TankRegulation
from app.models.regulations_master import RegulationsMaster
from app.models.cargo_tank import CargoTankTransaction
from app.models.cargo_master import CargoTankMaster
from app.models.tank_certificate import TankCertificate
from app.models.tank_drawings import TankDrawing
from app.models.valve_test_report import ValveTestReport

# --- CONFIGURATION ---
THEME_COLOR = RGBColor(0, 51, 102)
HEADER_TEXT = "Smart-Gas Pte Ltd"
SUB_HEADER_TEXT = "140 Paya Lebar Road #05-21, AZ Building, Singapore 409015\nTel : (65) 6848 1040 / Fax : (65) 6848 2173"

IMAGE_TYPE_MAP = {
    "frontview": "Front View", "rearview": "Rear View", "topview": "Top View",
    "undersideview": "Underside View", "frontlhview": "Front LH View",
    "rearlhview": "Rear LH View", "frontrhview": "Front RH View",
    "rearrhview": "Rear RH View", "lhsideview": "LH Side View",
    "rhsideview": "RH Side View", "valvessectionview": "Valves Section",
    "safetyvalve": "Safety Valve", "levelpressuregauge": "Level/Pressure Gauge",
    "vacuumreading": "Vacuum Reading"
}

def resolve_path(file_path, tank_number, base_dir):
    if not file_path: return None
    
    # Standardize slashes
    clean_db_path = str(file_path).replace("\\", "/")
    
    # Remove 'uploads/' prefix if it exists in the DB string
    if clean_db_path.startswith("uploads/"):
        clean_db_path = clean_db_path.replace("uploads/", "", 1)
        
    filename_only = os.path.basename(clean_db_path)

    # 1. The uploads folder where the code is currently running
    current_uploads = os.path.join(base_dir, "uploads")

    # 2. DYNAMIC RELATIVE PATH (No absolute "E:\" string)
    # This logic says: "Go up 2 levels, look for ISOTank-Mobile (1), then go into Backend/uploads"
    external_uploads = os.path.abspath(os.path.join(
        base_dir, 
        "..", "..",  # Go up to the drive/root folder
        "ISOTank-Mobile (1)", "Backend", "uploads" # Go down into the other project
    ))
    
    candidates = [
        # --- CHECK LOCAL FOLDER FIRST ---
        os.path.join(current_uploads, clean_db_path),
        os.path.join(current_uploads, "tank_images_mobile", clean_db_path),
        os.path.join(current_uploads, "tank_images_mobile", tank_number, "originals", filename_only),

        # --- CHECK NEIGHBOR PROJECT FOLDER ---
        os.path.join(external_uploads, clean_db_path),
        os.path.join(external_uploads, "tank_images_mobile", clean_db_path),
        os.path.join(external_uploads, "tank_images_mobile", tank_number, "originals", filename_only),
        os.path.join(external_uploads, "tank_images_mobile", tank_number, "thumbnails", filename_only),
        os.path.join(external_uploads, "drawings", tank_number, filename_only),
        os.path.join(external_uploads, "certificates", tank_number, filename_only)
    ]
    
    for path in candidates:
        if os.path.exists(path): return path
        
    return None
def format_value(value, suffix=""):
    if value is None or value == "": return "-"
    if isinstance(value, bool): return "Yes" if value else "No"
    if isinstance(value, (datetime, date)): return value.strftime("%d-%b-%Y")
    try:
        float_val = float(value)
        if float_val.is_integer(): return f"{int(float_val):,} {suffix}".strip()
        return f"{float_val:,.2f} {suffix}".strip()
    except: return f"{str(value)} {suffix}".strip()

# --- LAYOUT HELPERS (unchanged, simplified for brevity) ---
def add_custom_header(slide, title_text=""):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(0.5))
    p = txBox.text_frame.paragraphs[0]; p.text = HEADER_TEXT
    p.font.bold = True; p.font.size = Pt(24); p.font.color.rgb = THEME_COLOR; p.alignment = PP_ALIGN.CENTER
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.6), Inches(9), Inches(0.5))
    p = txBox.text_frame.paragraphs[0]; p.text = SUB_HEADER_TEXT
    p.font.size = Pt(9); p.font.color.rgb = RGBColor(80, 80, 80); p.alignment = PP_ALIGN.CENTER
    if title_text:
        shape = slide.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(9), Inches(0.4))
        shape.fill.solid(); shape.fill.fore_color.rgb = THEME_COLOR; shape.line.fill.background()
        shape.text_frame.text = title_text
        p = shape.text_frame.paragraphs[0]; p.font.bold = True; p.font.size = Pt(14); p.font.color.rgb = RGBColor(255, 255, 255); p.alignment = PP_ALIGN.LEFT

def create_section_title(slide, text_val, left, top, width=Inches(9)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.3))
    p = txBox.text_frame.paragraphs[0]; p.text = text_val; p.font.bold = True; p.font.size = Pt(11); p.font.color.rgb = THEME_COLOR
    return top + Inches(0.3)

def create_compact_table(slide, headers, data_rows, left, top, width, font_size=9):
    if not data_rows: return top
    rows = len(data_rows) + 1; cols = len(headers); row_height = Inches(0.25)
    table = slide.shapes.add_table(rows, cols, left, top, width, row_height * rows).table
    for i, h_text in enumerate(headers):
        cell = table.cell(0, i); cell.text = h_text; cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(230, 230, 230)
        p = cell.text_frame.paragraphs[0]; p.font.bold = True; p.font.size = Pt(font_size); p.font.color.rgb = RGBColor(0, 0, 0)
    for r_idx, row_data in enumerate(data_rows):
        for c_idx, value in enumerate(row_data):
            cell = table.cell(r_idx + 1, c_idx); cell.text = str(value)
            cell.text_frame.paragraphs[0].font.size = Pt(font_size)
    return top + (row_height * rows) + Inches(0.3)

def create_kv_block(slide, title, data_pairs, left, top, width):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.3))
    p = txBox.text_frame.paragraphs[0]; p.text = title; p.font.bold = True; p.font.size = Pt(11); p.font.color.rgb = THEME_COLOR
    table_top = top + Inches(0.3); rows = len(data_pairs); table_height = Inches(0.22) * rows
    table = slide.shapes.add_table(rows, 2, left, table_top, width, table_height).table
    table.columns[0].width = int(width * 0.45); table.columns[1].width = int(width * 0.55)
    for i, (label, value) in enumerate(data_pairs):
        cell_lbl = table.cell(i, 0); cell_lbl.text = str(label); cell_lbl.text_frame.paragraphs[0].font.bold = True; cell_lbl.text_frame.paragraphs[0].font.size = Pt(9)
        cell_val = table.cell(i, 1); cell_val.text = str(value); cell_val.text_frame.paragraphs[0].font.size = Pt(9)
    return table_top + table_height + Inches(0.2)

def add_image_sequence(prs, tank_number, image_list, section_title, base_dir):
    valid_items = []
    for item in image_list:
        if not item['path']: continue
        # PASS BASE_DIR HERE
        real_path = resolve_path(item['path'], tank_number, base_dir)
        if real_path: valid_items.append({'path': real_path, 'label': item['label']})
        else: valid_items.append({'path': None, 'label': item['label'], 'original': item['path']})

    if not valid_items: return
    
    for img_item in valid_items:
        path = img_item.get('path')
        label = img_item.get('label')
        current_slide = prs.slides.add_slide(prs.slide_layouts[5])
        add_custom_header(current_slide, f"{tank_number} - {section_title}")
        
        label_box = current_slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(9), Inches(0.4))
        label_box.text_frame.text = label
        p = label_box.text_frame.paragraphs[0]; p.font.bold = True; p.font.size = Pt(14); p.alignment = PP_ALIGN.CENTER
        
        img_top = Inches(2.2); img_left = Inches(0.5); max_height = Inches(5.0); max_width = Inches(9.0)

        if path:
            ext = os.path.splitext(path)[1].lower()
            if ext not in ['.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.gif', '.webp']:
                err_box = current_slide.shapes.add_textbox(img_left, img_top, max_width, Inches(1))
                err_box.text_frame.text = f"[File Format: {ext} - Cannot Render]"
                err_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            else:
                try:
                    pic = current_slide.shapes.add_picture(path, img_left, img_top, height=max_height)
                    if pic.width > max_width:
                        ratio = max_width / pic.width
                        pic.width = max_width
                        pic.height = int(pic.height * ratio)
                    pic.left = int((prs.slide_width - pic.width) / 2)
                except:
                    err_box = current_slide.shapes.add_textbox(img_left, img_top, max_width, Inches(1))
                    err_box.text_frame.text = "[Error Rendering Image]"
                    err_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    err_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)
        else:
            err_box = current_slide.shapes.add_textbox(img_left, img_top, max_width, Inches(1))
            err_box.text_frame.text = f"[Image Not Found]\n{img_item.get('original')}"
            err_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            err_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)

# ==============================================================================
# MAIN GENERATOR (Updated signature)
# ==============================================================================
def create_presentation(db: Session, tank_id: int, base_dir: str) -> BytesIO:
    tank = db.query(Tank).filter(Tank.id == tank_id).first()
    d = db.query(TankDetails).filter(TankDetails.tank_id == tank_id).first()
    if not tank or not d: raise ValueError("Tank Details not found")

    regs = db.query(TankRegulation, RegulationsMaster).outerjoin(RegulationsMaster, TankRegulation.regulation_id == RegulationsMaster.id).filter(TankRegulation.tank_id == tank_id).all()
    cargos = db.query(CargoTankTransaction, CargoTankMaster).join(CargoTankMaster, CargoTankTransaction.cargo_reference == CargoTankMaster.id).filter(CargoTankTransaction.tank_id == tank_id).all()
    certs = db.query(TankCertificate).filter(TankCertificate.tank_id == tank_id).all()
    drawings = db.query(TankDrawing).filter(TankDrawing.tank_id == tank_id).all()
    valves = db.query(ValveTestReport).filter(ValveTestReport.tank_id == tank_id).all()
    
    insp_sql = text("SELECT * FROM tank_inspection_details WHERE tank_id = :tid OR tank_number = :tn ORDER BY inspection_date DESC LIMIT 1")
    insp_row = db.execute(insp_sql, {"tid": tank_id, "tn": tank.tank_number}).mappings().first()
    insp = dict(insp_row) if insp_row else None

    checklist_rows = []
    todo_rows = []
    insp_images = []

    if insp:
        iid = insp['inspection_id']
        try:
            sql_chk = text("SELECT job_name, sub_job_description, status, comment FROM inspection_checklist WHERE inspection_id = :iid ORDER BY id ASC")
            checklist_data = db.execute(sql_chk, {"iid": iid}).fetchall()
            for r in checklist_data: checklist_rows.append([r[0], r[1], r[2] or "-", r[3] or "-"])
        except Exception: pass

        try:
            sql_todo = text("SELECT job_name, sub_job_description, status, comment FROM to_do_list WHERE inspection_id = :iid ORDER BY id ASC")
            todo_data = db.execute(sql_todo, {"iid": iid}).fetchall()
            for r in todo_data: todo_rows.append([r[0], r[1], r[2] or "Faulty", r[3] or "-"])
        except Exception: pass
            
        try:
            img_sql = text("SELECT image_type, image_path FROM tank_images WHERE inspection_id = :iid ORDER BY id ASC")
            img_data = db.execute(img_sql, {"iid": iid}).fetchall()
            for r in img_data:
                readable_label = IMAGE_TYPE_MAP.get(r[0], str(r[0]).title().replace("_", " "))
                insp_images.append({'path': r[1], 'label': readable_label})
        except Exception: pass
    
    if not insp_images:
        tank_images = db.query(TankImage).filter(TankImage.tank_number == tank.tank_number).all()
        for img in tank_images:
             readable_label = IMAGE_TYPE_MAP.get(img.image_type, (img.image_type or "").title())
             insp_images.append({'path': img.image_path, 'label': readable_label})

    # --- PPT GENERATION ---
    prs = Presentation()

    # SLIDE 1
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])
    add_custom_header(slide1, f"Executive Summary - {tank.tank_number}")
    gen_data = [("Manufacturer", d.mfgr), ("Date of Mfg", format_value(d.date_mfg)), ("Lease", format_value(d.lease)), ("UN ISO Code", d.un_iso_code), ("Size", d.size), ("Tare Weight", format_value(d.tare_weight_kg, "kg")), ("Gross Weight", format_value(d.gross_kg, "kg")), ("Max Payload", format_value(d.mpl_kg, "kg"))]
    create_kv_block(slide1, "General Identity", gen_data, Inches(0.5), Inches(1.8), Inches(4.2))

    if insp:
        insp_data = [("Report No", insp['report_number']), ("Inspection Date", format_value(insp['inspection_date'])), ("Status", insp['status_id']), ("Type", insp['inspection_type_id']), ("Inspector", insp['created_by']), ("Location", insp['location_id']), ("Safety Valve", f"{insp['safety_valve_brand_id'] or '-'} / {insp['safety_valve_model_id'] or '-'}"), ("Next Due", format_value(insp['pi_next_inspection_date']))]
        create_kv_block(slide1, "Latest Inspection", insp_data, Inches(5.0), Inches(1.8), Inches(4.5))
        if insp['notes']:
            txBox = slide1.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1)); txBox.text_frame.paragraphs[0].text = f"Inspector Notes: {insp['notes']}"; txBox.text_frame.paragraphs[0].font.size = Pt(10)

    # SLIDE 2
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    add_custom_header(slide2, f"Technical & Regulatory - {tank.tank_number}")
    tech_data = [("Capacity", format_value(d.capacity_l, "L")), ("MAWP", format_value(d.mawp, "Bar")), ("Working Pressure", format_value(d.working_pressure, "Bar")), ("Design Temp", d.design_temperature), ("Vessel Material", d.vesmat)]
    y_pos = create_kv_block(slide2, "Technical Specs", tech_data, Inches(0.5), Inches(1.8), Inches(4.2))
    const_data = [("Frame Type", d.frame_type), ("Cabinet", d.cabinet_type), ("Color", d.color_body_frame), ("Pump Type", d.pump_type)]
    create_kv_block(slide2, "Construction", const_data, Inches(0.5), y_pos, Inches(4.2))

    reg_rows = [[r.regulation_name if r else "Unknown", format_value(t.initial_approval_no), format_value(t.imo_type)] for t, r in regs]
    reg_y = create_section_title(slide2, "Regulations", Inches(5.0), Inches(1.8), Inches(4.5))
    create_compact_table(slide2, ["Regulation", "Approval", "IMO"], reg_rows, Inches(5.0), reg_y, Inches(4.5))

    # SLIDE 3
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    add_custom_header(slide3, f"Operational & Findings - {tank.tank_number}")
    current_y = Inches(1.8)
    
    cert_rows = [[c.certificate_number, format_value(c.inspection_agency), format_value(c.next_insp_date)] for c in certs]
    if cert_rows:
        current_y = create_section_title(slide3, "Certificates", Inches(0.5), current_y)
        current_y = create_compact_table(slide3, ["Cert No", "Agency", "Next Due"], cert_rows, Inches(0.5), current_y, Inches(9))

    cargo_rows = [[m.cargo_reference, format_value(c.density), format_value(c.loading_parts)] for c, m in cargos[:3]]
    if cargo_rows:
        current_y = create_section_title(slide3, "Approved Cargos (Recent)", Inches(0.5), current_y)
        current_y = create_compact_table(slide3, ["Cargo", "Density", "Parts"], cargo_rows, Inches(0.5), current_y, Inches(9))

    if todo_rows:
        current_y = create_section_title(slide3, "Faulty / To-Do Items", Inches(0.5), current_y)
        create_compact_table(slide3, ["Item", "Fault", "Status", "Comment"], todo_rows, Inches(0.5), current_y, Inches(9))
    elif checklist_rows:
        current_y = create_section_title(slide3, "Inspection Checklist (Preview)", Inches(0.5), current_y)
        create_compact_table(slide3, ["Category", "Check Item", "Status", "Comment"], checklist_rows[:8], Inches(0.5), current_y, Inches(9))

    # SLIDE 4+: IMAGES (Pass base_dir)
    add_image_sequence(prs, tank.tank_number, insp_images, "Inspection Photos", base_dir)
    
    cert_imgs = [{'path': c.certificate_file, 'label': f"Cert: {c.certificate_number}"} for c in certs if c.certificate_file]
    add_image_sequence(prs, tank.tank_number, cert_imgs, "Certificate Documents", base_dir)
    
    draw_imgs = [{'path': dr.file_path, 'label': f"Drawing: {dr.drawing_type}"} for dr in drawings if dr.file_path]
    add_image_sequence(prs, tank.tank_number, draw_imgs, "Technical Drawings", base_dir)
    
    valve_imgs = [{'path': v.inspection_report_file, 'label': f"Report: {format_value(v.test_date)}"} for v in valves if v.inspection_report_file]
    add_image_sequence(prs, tank.tank_number, valve_imgs, "Valve Test Documents", base_dir)

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output